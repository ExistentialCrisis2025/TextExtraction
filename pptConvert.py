import warnings
warnings.filterwarnings("ignore")

from pptx import Presentation
import os

def extract_ppt_content(ppt_path, output_dir="ppt_output"):
    prs = Presentation(ppt_path)
    os.makedirs(output_dir, exist_ok=True)

    all_content = {"slides": []}

    for i, slide in enumerate(prs.slides, start=1):
        slide_content = {"text": [], "tables": [], "images": []}

        # --- Extract text from shapes ---
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content["text"].append(shape.text.strip())

            # --- Extract tables ---
            if shape.has_table:
                table_data = []
                table = shape.table
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                slide_content["tables"].append(table_data)

            # --- Extract images ---
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                img_filename = os.path.join(output_dir, f"slide_{i}_image_{len(slide_content['images'])+1}.png")
                with open(img_filename, "wb") as f:
                    f.write(image_bytes)
                slide_content["images"].append(img_filename)

        all_content["slides"].append(slide_content)

    return all_content


if __name__ == "__main__":
    ppt_file = "Files/File_014.pptx"   # Change this to your ppt file
    content = extract_ppt_content(ppt_file)

    import json
    print(json.dumps(content, indent=2))
